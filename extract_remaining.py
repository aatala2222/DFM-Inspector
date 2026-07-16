"""Extract remaining large files with timeout protection"""
import os
from pathlib import Path

def extract_pdf_fast(filepath, max_pages=200):
    """Extract PDF with page limit and PyPDF2 (faster than pdfplumber)"""
    try:
        import PyPDF2
        text_parts = []
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            num_pages = min(len(reader.pages), max_pages)
            for i in range(num_pages):
                try:
                    text = reader.pages[i].extract_text()
                    if text:
                        text_parts.append(f"\n--- PAGE {i+1} ---\n{text}")
                except Exception:
                    continue
        return '\n'.join(text_parts)
    except Exception as e:
        return f"[ERROR: {e}]"


def extract_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n--- SLIDE {slide_num} ---")
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_parts.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                        if row_text.strip(' |'):
                            text_parts.append(row_text)
        return '\n'.join(text_parts)
    except Exception as e:
        return f"[ERROR: {e}]"


def extract_xlsx(filepath):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"\n--- SHEET: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip(' |'):
                    text_parts.append(row_text)
        return '\n'.join(text_parts)
    except Exception as e:
        return f"[ERROR: {e}]"


src = Path('Process Specs_043026')
out = Path('extracted_specs_new')
out.mkdir(exist_ok=True)

remaining = [
    'MachinerysHandbkOBERG.pdf',
    'Methods to Eliminate Surging Motion in a Conveyor System s12541-019-00042-y.pdf',
    'Plastics 101_Woojeon 2021.07.07 AR.pptx',
    'Sheet Metal Desin Guideline 930-00172_R01.pdf',
    'Sheet Metal Process Comparison Matrix.xlsx',
    'Springer Handbook of Mechanical Engineering 978-3-031-04329-1.pdf',
    'Stations BOM structure rev 02.pptx',
    'The Drum Motor 978-3-662-59298-4.pdf',
]

for fname in remaining:
    fp = src / fname
    if not fp.exists():
        print(f"SKIP (not found): {fname}")
        continue
    
    ext = fp.suffix.lower()
    out_name = fp.stem.replace(' ', '_').replace('(', '').replace(')', '') + '.txt'
    out_path = out / out_name
    
    if out_path.exists():
        print(f"SKIP (already done): {fname}")
        continue
    
    print(f"Processing: {fname}...")
    if ext == '.pdf':
        # Limit large PDFs
        max_pages = 100 if 'Handbk' in fname or 'Springer' in fname else 300
        content = extract_pdf_fast(fp, max_pages=max_pages)
    elif ext == '.pptx':
        content = extract_pptx(fp)
    elif ext == '.xlsx':
        content = extract_xlsx(fp)
    else:
        continue
    
    out_path.write_text(content, encoding='utf-8')
    print(f"  ✓ {len(content)/1024:.1f} KB")

print("\nDONE")
