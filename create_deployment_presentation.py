"""
Generate a PowerPoint presentation from the DFM Inspector deployment video script.
Creates a polished deck suitable for team training or stakeholder briefings.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# Color palette
COLOR_PRIMARY = RGBColor(0x1E, 0x40, 0xAF)      # Deep blue
COLOR_SECONDARY = RGBColor(0x3B, 0x82, 0xF6)    # Light blue
COLOR_ACCENT = RGBColor(0xF5, 0x9E, 0x0B)       # Amber
COLOR_SUCCESS = RGBColor(0x10, 0xB9, 0x81)      # Green
COLOR_DANGER = RGBColor(0xEF, 0x44, 0x44)       # Red
COLOR_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)       # Purple
COLOR_DARK = RGBColor(0x1E, 0x29, 0x3B)         # Slate-900
COLOR_GRAY = RGBColor(0x64, 0x74, 0x8B)         # Slate-500
COLOR_LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)   # Slate-100
COLOR_BG = RGBColor(0xF8, 0xFA, 0xFC)           # Slate-50
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_slide(prs, title, subtitle, presenter, role):
    """Title slide with branded layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Full-width colored band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.5))
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_PRIMARY
    band.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(12), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(12), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_DARK
    
    # Presenter
    pres_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(12), Inches(0.6))
    tf = pres_box.text_frame
    p = tf.paragraphs[0]
    p.text = presenter
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # Role
    role_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.0), Inches(12), Inches(0.5))
    tf = role_box.text_frame
    p = tf.paragraphs[0]
    p.text = role
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_GRAY
    p.font.italic = True


def add_section_divider(prs, section_num, section_name, color):
    """Section transition slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Full background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    
    # Big number
    num_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(12), Inches(2.5))
    tf = num_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_num
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Section name
    name_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(12), Inches(0.8))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_name
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE


def add_content_slide(prs, title, accent_color=COLOR_PRIMARY):
    """Returns slide and content area for further customization"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK
    
    # Title underline
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(2.5), Emu(40000))
    underline.fill.solid()
    underline.fill.fore_color.rgb = accent_color
    underline.line.fill.background()
    
    return slide


def add_bullet_text(slide, x, y, w, h, bullets, bullet_size=18, line_spacing=1.5):
    """Add bullet point text box"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(bullet_size)
        p.font.color.rgb = COLOR_DARK
        p.line_spacing = line_spacing
        p.space_after = Pt(8)


def add_textbox(slide, x, y, w, h, text, size=14, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    """Add a single text box"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return box


def add_card(slide, x, y, w, h, title, subtitle, color):
    """Add a colored card with title and subtitle"""
    # Background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    # Color accent on left
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.15), h)
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()
    
    # Title
    add_textbox(slide, x + Inches(0.3), y + Inches(0.15), w - Inches(0.4), Inches(0.5),
                title, size=16, bold=True, color=COLOR_DARK)
    
    # Subtitle
    add_textbox(slide, x + Inches(0.3), y + Inches(0.6), w - Inches(0.4), h - Inches(0.7),
                subtitle, size=12, color=COLOR_GRAY)


def add_code_block(slide, x, y, w, h, code, size=12):
    """Add a monospace code block on a dark background"""
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_DARK
    bg.line.fill.background()
    
    box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), h - Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = True
    
    lines = code.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.name = 'Consolas'
        p.font.color.rgb = RGBColor(0xE0, 0xF2, 0xFE)


def main():
    prs = Presentation()
    
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ---------- SLIDE 1: TITLE ----------
    add_title_slide(prs,
                    "DFM Inspector",
                    "From Local Python App to Team-Shared Dashboard",
                    "Alex Atala  |  Senior Product Development Engineer",
                    "Amazon Robotics")
    
    # ---------- SLIDE 2: AGENDA ----------
    slide = add_content_slide(prs, "What We'll Cover")
    
    items = [
        ("Phase 1", "Local Development", COLOR_SECONDARY),
        ("Phase 2", "Source Control Backup", COLOR_ACCENT),
        ("Phase 3", "Code Review & Merge", COLOR_DANGER),
        ("Phase 4", "Sharing With Team", COLOR_PURPLE),
    ]
    
    y = Inches(1.6)
    for label, name, color in items:
        # Phase number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(1.0), y, Inches(2.2), Inches(1.0))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        
        add_textbox(slide, Inches(1.0), y + Inches(0.25), Inches(2.2), Inches(0.5),
                    label, size=22, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        
        # Phase description
        add_textbox(slide, Inches(3.5), y + Inches(0.25), Inches(8), Inches(0.6),
                    name, size=24, bold=True, color=COLOR_DARK)
        
        y += Inches(1.3)
    
    # ---------- SLIDE 3: END GOAL ----------
    slide = add_content_slide(prs, "End Goal", accent_color=COLOR_SUCCESS)
    
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.8),
                "A working web app where engineers upload STEP files",
                size=24, color=COLOR_DARK)
    add_textbox(slide, Inches(0.7), Inches(2.1), Inches(12), Inches(0.8),
                "and get DFM analysis based on Amazon Robotics design guidelines.",
                size=24, color=COLOR_DARK)
    
    # 3 outcome cards
    add_card(slide, Inches(0.7), Inches(3.5), Inches(4), Inches(2.8),
             "📋 Standards-Based",
             "Rules from 930-00172 (sheet metal), 930-00166 (casting), NADCA, ISO 2768, AWS",
             COLOR_PRIMARY)
    
    add_card(slide, Inches(4.9), Inches(3.5), Inches(4), Inches(2.8),
             "🔍 Real Geometry",
             "CadQuery + OpenCascade B-Rep extraction detects holes, bends, thickness, spatial relationships",
             COLOR_ACCENT)
    
    add_card(slide, Inches(9.1), Inches(3.5), Inches(4), Inches(2.8),
             "🌐 Shareable",
             "Local install for power users + AgentSpaces dashboard for everyone else",
             COLOR_SUCCESS)
    
    # =====================================================================
    # PHASE 1: LOCAL DEVELOPMENT
    # =====================================================================
    add_section_divider(prs, "1", "Local Development", COLOR_SECONDARY)
    
    # Slide: Phase 1 Overview
    slide = add_content_slide(prs, "Phase 1: Local Development", COLOR_SECONDARY)
    
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.6),
                "Goal: Get the app running on your local machine",
                size=18, color=COLOR_GRAY, italic=True)
    
    add_textbox(slide, Inches(0.7), Inches(2.4), Inches(12), Inches(0.5),
                "Three Steps:",
                size=20, bold=True, color=COLOR_DARK)
    
    add_card(slide, Inches(0.7), Inches(3.0), Inches(4), Inches(3.5),
             "1.1 Install Dependencies",
             "Python 3.11+ required. Run pip install -r requirements.txt to get Flask, Trimesh, Cascadio, NumPy, CadQuery, and python-docx.",
             COLOR_SECONDARY)
    
    add_card(slide, Inches(4.9), Inches(3.0), Inches(4), Inches(3.5),
             "1.2 Test Locally",
             "Run python app.py — Flask server starts on port 5000. Visit http://localhost:5000 in your browser. Upload a STEP file to verify analysis.",
             COLOR_SECONDARY)
    
    add_card(slide, Inches(9.1), Inches(3.0), Inches(4), Inches(3.5),
             "1.3 Desktop Shortcut",
             "Create one-click launcher with a .bat file + PowerShell script. Double-click to start app + open browser.",
             COLOR_SECONDARY)
    
    # Slide: Phase 1 Commands
    slide = add_content_slide(prs, "Phase 1: Commands", COLOR_SECONDARY)
    
    add_textbox(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
                "Install & Run", size=18, bold=True, color=COLOR_DARK)
    
    add_code_block(slide, Inches(0.7), Inches(2.0), Inches(12), Inches(1.2),
                   "pip install -r requirements.txt\npython app.py")
    
    add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12), Inches(0.5),
                "Verify in browser", size=18, bold=True, color=COLOR_DARK)
    
    add_code_block(slide, Inches(0.7), Inches(4.0), Inches(12), Inches(0.8),
                   "http://localhost:5000")
    
    # Warning
    warn_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.7), Inches(5.2), Inches(12), Inches(1.0))
    warn_box.fill.solid()
    warn_box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    warn_box.line.color.rgb = COLOR_ACCENT
    warn_box.line.width = Pt(1.5)
    
    add_textbox(slide, Inches(0.9), Inches(5.35), Inches(11.6), Inches(0.5),
                "⚠ Common gotcha: Disconnect Zscaler/VPN when testing localhost",
                size=14, bold=True, color=COLOR_DARK)
    add_textbox(slide, Inches(0.9), Inches(5.7), Inches(11.6), Inches(0.5),
                "Some VPN clients intercept localhost traffic — the local app doesn't need VPN.",
                size=12, color=COLOR_GRAY)
    
    # =====================================================================
    # PHASE 2: SOURCE CONTROL
    # =====================================================================
    add_section_divider(prs, "2", "Source Control Backup", COLOR_ACCENT)
    
    # Slide: Phase 2 Overview
    slide = add_content_slide(prs, "Phase 2: Backup to code.amazon.com", COLOR_ACCENT)
    
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.6),
                "Goal: Get your code into Amazon source control",
                size=18, color=COLOR_GRAY, italic=True)
    
    steps = [
        ("2.1", "Connect to Amazon VPN", "Cisco AnyConnect — required for git.amazon.com access"),
        ("2.2", "Authenticate with Midway", "mwinit -f generates SSH cert for ~20 hours"),
        ("2.3", "Create Brazil Package", "code.amazon.com → Create Package → Python Brazil template"),
        ("2.4", "Clone the Empty Package", "git clone ssh://git.amazon.com:2222/pkg/<PackageName>"),
        ("2.5", "Copy Code In", "src/<package_name>/ for modules; root for app.py, etc."),
        ("2.6", "Push to Feature Branch", "git push origin HEAD:refs/heads/initial-import"),
    ]
    
    y = Inches(2.4)
    for num, name, desc in steps:
        # Step number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(0.7), y, Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_ACCENT
        circle.line.fill.background()
        
        add_textbox(slide, Inches(0.7), y + Inches(0.1), Inches(0.7), Inches(0.5),
                    num, size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        
        # Step name
        add_textbox(slide, Inches(1.6), y + Inches(0.05), Inches(4), Inches(0.4),
                    name, size=14, bold=True, color=COLOR_DARK)
        
        # Description
        add_textbox(slide, Inches(1.6), y + Inches(0.4), Inches(11), Inches(0.4),
                    desc, size=11, color=COLOR_GRAY)
        
        y += Inches(0.75)
    
    # Slide: Phase 2 Critical Gotcha
    slide = add_content_slide(prs, "Phase 2: Critical Gotcha", COLOR_ACCENT)
    
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.7),
                "Don't push directly to mainline",
                size=24, bold=True, color=COLOR_DANGER)
    
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12), Inches(0.6),
                "GitFarm protects mainline — direct pushes are rejected.",
                size=16, color=COLOR_DARK)
    
    add_code_block(slide, Inches(0.7), Inches(3.0), Inches(12), Inches(1.2),
                   "# ❌ Wrong - will fail with hook decline\ngit push origin HEAD:mainline")
    
    add_code_block(slide, Inches(0.7), Inches(4.4), Inches(12), Inches(1.2),
                   "# ✓ Correct - push to feature branch\ngit push origin HEAD:refs/heads/initial-import")
    
    add_textbox(slide, Inches(0.7), Inches(5.9), Inches(12), Inches(0.6),
                "Then create a Code Review (CR) to merge into mainline — covered in Phase 3.",
                size=14, color=COLOR_GRAY, italic=True)
    
    # Slide: Phase 2 Folder Structure Gotcha
    slide = add_content_slide(prs, "Phase 2: Watch the Folder Names", COLOR_ACCENT)
    
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.6),
                "Windows is case-insensitive — folder names can collide with Brazil files",
                size=16, color=COLOR_DARK)
    
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(0.5),
                "❌ Don't do this:", size=16, bold=True, color=COLOR_DANGER)
    
    add_code_block(slide, Inches(0.7), Inches(2.9), Inches(5.5), Inches(2.5),
                   "<Package>/\n├── Config       (Brazil)\n├── config/      (your YAML)\n│   └── ...\n└── ...\n\n# Conflicts on Windows!")
    
    add_textbox(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(0.5),
                "✓ Do this instead:", size=16, bold=True, color=COLOR_SUCCESS)
    
    add_code_block(slide, Inches(7), Inches(2.9), Inches(5.5), Inches(2.5),
                   "<Package>/\n├── Config       (Brazil)\n├── dfm_config/  (your YAML)\n│   └── ...\n└── ...\n\n# No conflict")
    
    add_textbox(slide, Inches(0.7), Inches(5.7), Inches(12), Inches(0.6),
                "Same applies to test/ vs tests/ — Brazil uses singular 'test'",
                size=14, color=COLOR_GRAY, italic=True)
    
    # =====================================================================
    # PHASE 3: MERGE BOTTLENECK
    # =====================================================================
    add_section_divider(prs, "3", "Code Review & Merge", COLOR_DANGER)
    
    # Slide: Phase 3 The Bottleneck
    slide = add_content_slide(prs, "Phase 3: The Bottleneck", COLOR_DANGER)
    
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.7),
                "Merging to mainline requires the cr CLI tool",
                size=22, bold=True, color=COLOR_DARK)
    
    add_textbox(slide, Inches(0.7), Inches(2.4), Inches(12), Inches(0.6),
                "cr is NOT available on Windows. The internal Amazon repos require:",
                size=16, color=COLOR_DARK)
    
    bullets = [
        "Cloud Desktop access (recommended), OR",
        "WSL with proper internal network access, OR",
        "A Linux dev environment with toolbox installed",
    ]
    add_bullet_text(slide, Inches(1.3), Inches(3.2), Inches(11), Inches(2),
                    bullets, bullet_size=15)
    
    # Failed attempts box
    fail_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.7), Inches(5.0), Inches(12), Inches(1.8))
    fail_box.fill.solid()
    fail_box.fill.fore_color.rgb = RGBColor(0xFE, 0xE2, 0xE2)
    fail_box.line.color.rgb = COLOR_DANGER
    fail_box.line.width = Pt(1.5)
    
    add_textbox(slide, Inches(0.9), Inches(5.15), Inches(11.6), Inches(0.5),
                "What I tried (none worked from Windows):",
                size=14, bold=True, color=COLOR_DARK)
    
    add_textbox(slide, Inches(0.9), Inches(5.55), Inches(11.6), Inches(1.3),
                "• toolbox install cr  → command not found    • pip install cruxcr  → not on public PyPI\n"
                "• pip install amzn-cr  → internal artifactory unreachable    • Web 'Create Review' button doesn't exist\n"
                "• Direct push to mainline  → blocked by GitFarm branch protection",
                size=12, color=COLOR_GRAY)
    
    # Slide: Phase 3 Three Options
    slide = add_content_slide(prs, "Phase 3: Three Paths Forward", COLOR_DANGER)
    
    add_card(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(5.3),
             "Option A: Cloud Desktop",
             "RECOMMENDED LONG-TERM\n\n"
             "Request via clouddesktop.amazon.com\n\n"
             "Pre-installed tools:\n• cr CLI\n• brazil\n• toolbox\n• mwinit\n\n"
             "30 seconds to create CR once provisioned.\n\nEffort: Manager approval + 1-2 days",
             COLOR_DANGER)
    
    add_card(slide, Inches(4.7), Inches(1.5), Inches(4), Inches(5.3),
             "Option B: Ask Teammate",
             "FASTEST\n\n"
             "Send teammate with Cloud Desktop:\n• Package URL\n• Branch name (initial-import)\n• Reviewer alias\n\n"
             "They run cr from their machine — 2 min task.\n\nEffort: Slack DM + their availability",
             COLOR_DANGER)
    
    add_card(slide, Inches(8.9), Inches(1.5), Inches(4), Inches(5.3),
             "Option C: Defer",
             "SHIP NOW\n\n"
             "Leave code on initial-import branch.\n\nCode is safely backed up.\n\n"
             "You CAN still:\n• Develop locally\n• Push updates\n• Build AgentSpaces dashboard\n\n"
             "You CAN'T:\n• Set up pipeline\n• Deploy via Apollo",
             COLOR_DANGER)
    
    # =====================================================================
    # PHASE 4: SHARING
    # =====================================================================
    add_section_divider(prs, "4", "Share With Team", COLOR_PURPLE)
    
    # Slide: Phase 4 Path A vs Path B
    slide = add_content_slide(prs, "Phase 4: Two Sharing Paths", COLOR_PURPLE)
    
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.6),
                "Pick based on your team's needs — or do both",
                size=18, color=COLOR_GRAY, italic=True)
    
    # Path A
    path_a = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.7), Inches(2.4), Inches(6), Inches(4.5))
    path_a.fill.solid()
    path_a.fill.fore_color.rgb = RGBColor(0xFA, 0xF5, 0xFF)
    path_a.line.color.rgb = COLOR_PURPLE
    path_a.line.width = Pt(2)
    
    add_textbox(slide, Inches(0.9), Inches(2.6), Inches(5.6), Inches(0.6),
                "📚 Path A: Documentation",
                size=22, bold=True, color=COLOR_PURPLE)
    
    add_textbox(slide, Inches(0.9), Inches(3.3), Inches(5.6), Inches(0.5),
                "AgentSpaces HTML Dashboard",
                size=16, bold=True, color=COLOR_DARK)
    
    add_bullet_text(slide, Inches(0.9), Inches(3.8), Inches(5.6), Inches(2.5),
                    ["Upload DFM_RULES_REFERENCE.md",
                     "Add source PDFs from Process Specs",
                     "Ask AI: \"Generate HTML dashboard\"",
                     "Get shareable URL",
                     "Send to team via Slack"], bullet_size=12)
    
    add_textbox(slide, Inches(0.9), Inches(6.3), Inches(5.6), Inches(0.5),
                "⏱ Effort: ~1 hour",
                size=12, color=COLOR_GRAY, italic=True)
    
    # Path B
    path_b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(7), Inches(2.4), Inches(6), Inches(4.5))
    path_b.fill.solid()
    path_b.fill.fore_color.rgb = RGBColor(0xFA, 0xF5, 0xFF)
    path_b.line.color.rgb = COLOR_PURPLE
    path_b.line.width = Pt(2)
    
    add_textbox(slide, Inches(7.2), Inches(2.6), Inches(5.6), Inches(0.6),
                "⚙ Path B: Live App",
                size=22, bold=True, color=COLOR_PURPLE)
    
    add_textbox(slide, Inches(7.2), Inches(3.3), Inches(5.6), Inches(0.5),
                "Working DFM Analyzer for Team",
                size=16, bold=True, color=COLOR_DARK)
    
    add_bullet_text(slide, Inches(7.2), Inches(3.8), Inches(5.6), Inches(2.5),
                    ["B1: Local installer (1-2 days)",
                     "B2: Shared Cloud Desktop (3-5 days)",
                     "B3: Pipeline + Apollo (1-3 weeks)",
                     "Pick based on team size + criticality"],
                    bullet_size=12)
    
    add_textbox(slide, Inches(7.2), Inches(6.3), Inches(5.6), Inches(0.5),
                "⏱ Effort: Days to weeks",
                size=12, color=COLOR_GRAY, italic=True)
    
    # Slide: Phase 4 Path B Methods
    slide = add_content_slide(prs, "Phase 4 Path B: Three Deployment Methods", COLOR_PURPLE)
    
    methods = [
        ("B1", "Local Installer", "1-2 days",
         "Each user installs locally\nDistribute via Quip/SharePoint\nSetup script (PowerShell)\nDesktop shortcut auto-created",
         "Small teams (5-10 users)", COLOR_SUCCESS),
        ("B2", "Shared Cloud Desktop", "3-5 days",
         "Run app on team Cloud Desktop\nPort 5000 open to corp network\nProcess supervisor for uptime\nSingle URL for everyone",
         "Medium teams (10-30 users)", COLOR_ACCENT),
        ("B3", "Pipeline + Apollo", "1-3 weeks",
         "Production WSGI (gunicorn)\nMidway authentication\nS3 storage, CloudWatch logs\nAppSec review required",
         "Large teams (50+ users)", COLOR_DANGER),
    ]
    
    x = Inches(0.5)
    for code, name, effort, details, audience, badge_color in methods:
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, Inches(1.5), Inches(4.1), Inches(5.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        card.line.color.rgb = COLOR_PURPLE
        card.line.width = Pt(1.5)
        
        # Method code badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       x + Inches(0.3), Inches(1.7), Inches(0.7), Inches(0.7))
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_PURPLE
        badge.line.fill.background()
        add_textbox(slide, x + Inches(0.3), Inches(1.78), Inches(0.7), Inches(0.5),
                    code, size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        
        # Name
        add_textbox(slide, x + Inches(1.1), Inches(1.7), Inches(2.9), Inches(0.5),
                    name, size=15, bold=True, color=COLOR_DARK)
        
        # Effort
        add_textbox(slide, x + Inches(1.1), Inches(2.05), Inches(2.9), Inches(0.4),
                    f"⏱ {effort}", size=11, color=COLOR_GRAY, italic=True)
        
        # Details
        add_textbox(slide, x + Inches(0.3), Inches(2.7), Inches(3.7), Inches(2.8),
                    details, size=11, color=COLOR_DARK)
        
        # Audience badge
        aud_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x + Inches(0.3), Inches(6.0), Inches(3.5), Inches(0.7))
        aud_box.fill.solid()
        aud_box.fill.fore_color.rgb = badge_color
        aud_box.line.fill.background()
        
        add_textbox(slide, x + Inches(0.3), Inches(6.15), Inches(3.5), Inches(0.5),
                    audience, size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        
        x += Inches(4.3)
    
    # =====================================================================
    # WRAP-UP
    # =====================================================================
    
    # Slide: Recap
    slide = add_content_slide(prs, "Recap")
    
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.6),
                "From local Python app to team-shared dashboard:",
                size=18, color=COLOR_GRAY, italic=True)
    
    recap_items = [
        ("1", "Local Development", "Hours", "Get app running with desktop shortcut", COLOR_SECONDARY),
        ("2", "Source Control Backup", "Hour", "Push to code.amazon.com feature branch", COLOR_ACCENT),
        ("3", "Code Review & Merge", "Day", "Use cr CLI on Cloud Desktop", COLOR_DANGER),
        ("4A", "AgentSpaces Dashboard", "Hour", "Generate shareable rules HTML", COLOR_PURPLE),
        ("4B", "Live App Deployment", "Days-Weeks", "Local install / Cloud Desktop / Pipeline", COLOR_PURPLE),
    ]
    
    y = Inches(2.4)
    for num, name, time, desc, color in recap_items:
        # Number
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(0.7), y, Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_textbox(slide, Inches(0.7), y + Inches(0.1), Inches(0.7), Inches(0.5),
                    num, size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        
        # Name
        add_textbox(slide, Inches(1.6), y, Inches(4.5), Inches(0.5),
                    name, size=15, bold=True, color=COLOR_DARK)
        
        # Time pill
        time_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(6.0), y + Inches(0.05), Inches(1.5), Inches(0.5))
        time_box.fill.solid()
        time_box.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        time_box.line.color.rgb = COLOR_GRAY
        time_box.line.width = Pt(0.75)
        add_textbox(slide, Inches(6.0), y + Inches(0.13), Inches(1.5), Inches(0.4),
                    time, size=11, bold=True, color=COLOR_GRAY, align=PP_ALIGN.CENTER)
        
        # Description
        add_textbox(slide, Inches(7.7), y + Inches(0.1), Inches(5.3), Inches(0.5),
                    desc, size=12, color=COLOR_GRAY)
        
        y += Inches(0.85)
    
    # Slide: Quick Win
    slide = add_content_slide(prs, "Minimum Viable Path", COLOR_SUCCESS)
    
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.7),
                "If you only have 1 day, do this:",
                size=22, bold=True, color=COLOR_DARK)
    
    add_card(slide, Inches(1.5), Inches(2.7), Inches(4.5), Inches(3.2),
             "✓ Phase 1: Local App",
             "Get app working locally\nDesktop shortcut\n~2 hours",
             COLOR_SECONDARY)
    
    # Plus
    add_textbox(slide, Inches(6.0), Inches(3.7), Inches(1.5), Inches(1.5),
                "+", size=72, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)
    
    add_card(slide, Inches(7.5), Inches(2.7), Inches(4.5), Inches(3.2),
             "✓ Phase 4A: AgentSpaces Dashboard",
             "Upload DFM rules MD\nAsk AI for HTML dashboard\nShare URL with team\n~1 hour",
             COLOR_PURPLE)
    
    # Result
    result = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.7), Inches(6.2), Inches(12), Inches(0.9))
    result.fill.solid()
    result.fill.fore_color.rgb = COLOR_SUCCESS
    result.line.fill.background()
    
    add_textbox(slide, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
                "= Personal working app + Shared rules reference for the team",
                size=18, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    
    # Slide: Resources
    slide = add_content_slide(prs, "Resources")
    
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.5),
                "All deliverables in C:\\Users\\<you>\\DFM Inspector\\",
                size=14, color=COLOR_GRAY, italic=True)
    
    docs = [
        ("DEPLOYMENT_GUIDE.md", "Full step-by-step walkthrough"),
        ("DFM_Deployment_Checklist.md", "Printable one-page checklist"),
        ("DFM_Deployment_Flowchart.svg", "Visual process diagram"),
        ("DFM_Deployment_Video_Script.md", "Screencast narration script"),
        ("DFM_RULES_REFERENCE.md", "Catalog of all DFM rules"),
        ("DFM_Inspector_Launcher.bat", "Local app launcher"),
        ("create_dfm_violation_step.py", "Test part generator"),
    ]
    
    y = Inches(2.3)
    for fname, desc in docs:
        # Icon
        icon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.9), y, Inches(0.4), Inches(0.5))
        icon.fill.solid()
        icon.fill.fore_color.rgb = COLOR_PRIMARY
        icon.line.fill.background()
        add_textbox(slide, Inches(0.9), y + Inches(0.05), Inches(0.4), Inches(0.4),
                    "📄", size=18, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        
        # Filename
        add_textbox(slide, Inches(1.5), y, Inches(5), Inches(0.5),
                    fname, size=13, bold=True, color=COLOR_DARK)
        
        # Description
        add_textbox(slide, Inches(6.5), y, Inches(6.5), Inches(0.5),
                    desc, size=12, color=COLOR_GRAY)
        
        y += Inches(0.6)
    
    # Slide: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()
    
    add_textbox(slide, Inches(0.7), Inches(2.5), Inches(12), Inches(1.5),
                "Thanks for watching",
                size=60, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(0.7), Inches(4.0), Inches(12), Inches(0.7),
                "Now go ship something.",
                size=24, italic=True, color=COLOR_LIGHT_GRAY, align=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(0.7), Inches(6.0), Inches(12), Inches(0.5),
                "Alex Atala  |  Senior Product Development Engineer",
                size=14, color=COLOR_LIGHT_GRAY, align=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
                "Amazon Robotics",
                size=12, italic=True, color=COLOR_LIGHT_GRAY, align=PP_ALIGN.CENTER)
    
    # Save
    output = 'DFM_Inspector_Deployment_Presentation.pptx'
    prs.save(output)
    print(f"Created: {output}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
