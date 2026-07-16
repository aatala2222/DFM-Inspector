"""
Extract text content from all documents in Process Specs_043026/
for DFM rule analysis and analyzer updates.
"""
import os
import sys
from pathlib import Path


def extract_pdf(filepath):
    """Extract text from PDF using pdfplumber"""
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(filepath) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ''
                text_parts.append(f"\n--- PAGE {page_num} ---\n{text}")
        return '\n'.join(text_parts)
    except Exception as e:
        return f"[ERROR extracting PDF: {e}]"


def extract_docx(filepath):
    """Extract text from DOCX"""
    try:
        from docx import Document
        doc = Document(filepath)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text.strip(' |'):
                    text_parts.append(row_text)
        return '\n'.join(text_parts)
    except Exception as e:
        return f"[ERROR extracting DOCX: {e}]"


def extract_pptx(filepath):
    """Extract text from PPTX"""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n--- SLIDE {slide_num} ---")
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_parts.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                        if row_text.strip(' |'):
                            text_parts.append(row_text)
        return '\n'.join(text_parts)
    except Exception as e:
        return f"[ERROR extracting PPTX: {e}]"


def extract_xlsx(filepath):
    """Extract text from XLSX"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"\n--- SHEET: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip(' |'):
                    text_parts.append(row_text)
        return '\n'.join(text_parts)
    except Exception as e:
        return f"[ERROR extracting XLSX: {e}]"


def main():
    src_dir = Path('Process Specs_043026')
    output_dir = Path('extracted_specs_new')
    output_dir.mkdir(exist_ok=True)
    
    if not src_dir.exists():
        print(f"Source directory not found: {src_dir}")
        return
    
    summary = []
    
    for filepath in sorted(src_dir.iterdir()):
        if filepath.is_file() and not filepath.name.startswith('.'):
            ext = filepath.suffix.lower()
            print(f"Processing: {filepath.name}...")
            
            content = None
            if ext == '.pdf':
                content = extract_pdf(filepath)
            elif ext == '.docx':
                content = extract_docx(filepath)
            elif ext == '.pptx':
                content = extract_pptx(filepath)
            elif ext == '.xlsx':
                content = extract_xlsx(filepath)
            elif ext in ('.png', '.jpg', '.jpeg'):
                print(f"  Skipping image file: {filepath.name}")
                continue
            else:
                print(f"  Unsupported format: {ext}")
                continue
            
            if content:
                out_name = filepath.stem.replace(' ', '_').replace('(', '').replace(')', '') + '.txt'
                out_path = output_dir / out_name
                out_path.write_text(content, encoding='utf-8')
                size_kb = len(content) / 1024
                summary.append(f"  {filepath.name} → {out_name} ({size_kb:.1f} KB)")
                print(f"  ✓ Extracted {size_kb:.1f} KB")
    
    print("\n" + "=" * 60)
    print("EXTRACTION COMPLETE")
    print("=" * 60)
    for line in summary:
        print(line)


if __name__ == '__main__':
    main()
